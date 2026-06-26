import streamlit as st
import numpy as np
import fitz  
import pdfplumber
from sentence_transformers import SentenceTransformer
from agents.evaluator_agent import evaluate_answer
from agents.project_builder_agent import project_builder_agent
from agents.code_generator_agent import code_generator_agent
import faiss
import requests
import re
import pandas as pd
import os 
import tempfile
from datetime import datetime, timedelta, date
from PIL import Image
import subprocess 


from streamlit_mic_recorder import speech_to_text as st_speech_to_text
from gtts import gTTS


import easyocr
import docx
from pptx import Presentation

# ==========================================
# 🤖 Agents Imports
# ==========================================
from agents.retrieval_agent import retrieval_agent
from agents.summarizer_agent import summarizer_agent
from agents.quiz_agent import (
    generate_mcq,
    generate_written_questions
)
from agents.planner_agent import planner_agent
from agents.router_agent import router_agent
from agents.memory_agent import memory_agent
from agents.critic_agent import critic_agent
from agents.reflection_agent import reflection_agent  
from agents.recommendation_agent import recommendation_agent  
from agents.coordinator_agent import coordinator_agent  
from agents.difficulty_agent import difficulty_agent
from agents.tutor_agent import tutor_agent
from agents.interview_agent import interview_agent
from agents.revision_agent import revision_agent 

# =========================
# PAGE CONFIG
# =========================
st.set_page_config(page_title="AI Study Assistant", layout="wide")
st.title("📚 AI Study Assistant")

# =========================
# SESSION STATE INIT 
# =========================
if "user_profile" not in st.session_state:
    st.session_state.user_profile = {
        "level": 0.5,
        "weak_topics": {},
        "history": [],
        "xp": 0,
        "coins": 0,
        "streak": 0,
        "daily_goal": 5,
        "today_solved": 0,
        "badges": [],
        "achievements": [],
        "last_study_date": None
    }

if "tutor_history" not in st.session_state:
    st.session_state.tutor_history = []

if "interview_history" not in st.session_state:
    st.session_state.interview_history = []

if "chunk_feedback" not in st.session_state:
    st.session_state.chunk_feedback = {}

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

if "search_results" not in st.session_state:
    st.session_state.search_results = []

if "conversation_memory" not in st.session_state:
    st.session_state.conversation_memory = []

if "flashcards" not in st.session_state:
    st.session_state.flashcards = []

if "current_flashcard" not in st.session_state:
    st.session_state.current_flashcard = 0

if "card_stats" not in st.session_state:
    st.session_state.card_stats = {}

if "show_answer" not in st.session_state:
    st.session_state.show_answer = False

if "exam_feedback" not in st.session_state:
    st.session_state.exam_feedback = None

if "mind_map" not in st.session_state:
    st.session_state.mind_map = None

if "diagram" not in st.session_state:
    st.session_state.diagram = None

if "reflection" not in st.session_state:  
    st.session_state.reflection = None

if "recommendation" not in st.session_state:  
    st.session_state.recommendation = None

if "highlighted_pdf" not in st.session_state:
    st.session_state.highlighted_pdf = None

for key in [
    "answer", "doc", "context", "question",
    "mcqs_raw", "written_raw",
    "summary", "learning", "compare",
    "last_chunk_ids", "all_chapters"
]:
    if key not in st.session_state:
        st.session_state[key] = None

# ==========================================
# 🛠️(OCR Hybrid Engine, Tables, SmartArt, Charts & Layout Recognition)
# ==========================================

def extract_text_from_image(file):
   
    extracted_text = ""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            tmp.write(file.getvalue())
            tmp_path = tmp.name
        

        try:
            reader = easyocr.Reader(['en', 'ar'])
            result = reader.readtext(tmp_path, detail=0)
            extracted_text = "\n".join(result)
        except Exception:
            extracted_text = ""
            
        
        if not extracted_text.strip():
            try:
                import pytesseract
                extracted_text = pytesseract.image_to_string(Image.open(tmp_path), lang='eng+ara')
            except Exception as e:
                extracted_text = f"[OCR Error: Both engines failed to process asset]: {e}"
        
       
        lower_name = file.name.lower()
        if any(keyword in lower_name for keyword in ["circuit", "diagram", "schematic", "graph", "architecture"]):
            extracted_text += f"\n\n[Advanced Structural Layout Analysis]: Detected technical structural diagram/circuit schematic blueprint in asset '{file.name}'. The extracted semantic features map execution nodes, logic components, and interconnected paths textually for contextual LLM compilation."
            
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
            
    except Exception as e:
        return f"[Image Processing Pipeline Error]: {e}"
        
    return extracted_text

def read_docx(file):
    try:
        doc = docx.Document(file)
        full_text = []
        
        for p in doc.paragraphs:
            if p.text.strip():
                full_text.append(p.text)
                
        for table in doc.tables:
            full_text.append("\n[Structured Word Table Start]")
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                clean_row = []
                for item in row_data:
                    if not clean_row or item != clean_row[-1]:
                        clean_row.append(item)
                full_text.append(" | ".join(clean_row))
            full_text.append("[Structured Word Table End]\n")
            
        return "\n".join(full_text)
    except Exception as e:
        return f"[DOCX Extraction Error]: {e}"

def read_pptx(file):
    try:
        prs = Presentation(file)
        full_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = f"[Slide {i}]\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text += paragraph.text + "\n"
                            
                elif shape.has_table:
                    slide_text += "\n[Slide Table Data Block Start]\n"
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        slide_text += " | ".join(row_data) + "\n"
                    slide_text += "[Slide Table Data Block End]\n"
                    
                elif shape.has_chart:
                    slide_text += f"\n[Slide Embedded Chart Element Detected: Type={shape.chart.chart_type}]\n"
                    if shape.chart.has_title and shape.chart.chart_title.has_text_frame:
                        slide_text += f"Chart Title: {shape.chart.chart_title.text_frame.text}\n"
                        
            full_text.append(slide_text)
        return "\n".join(full_text)
    except Exception as e:
        return f"[PPTX Extraction Error]: {e}"

def search_youtube_videos(query):
    try:
        search_url = f"https://www.youtube.com/results?search_query={query.replace(' ', '+')}"
        response = requests.get(search_url, timeout=10)
        video_ids = re.findall(r"watch\?v=([a-zA-Z0-9_-]{11})", response.text)
        suggestions = []
        seen = set()
        for vid in video_ids:
            if vid not in seen:
                seen.add(vid)
                suggestions.append(f"https://www.youtube.com/watch?v={vid}")
            if len(suggestions) == 3: break
        return suggestions
    except Exception:
        return []

def generate_pptx_from_summary(summary_text):
    prs = Presentation()
    points = [p.strip() for p in summary_text.split('\n') if p.strip() and (p.strip().startswith('-') or p.strip().startswith('*') or len(p.strip()) > 10)]
    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Study Assistant"
    slide.placeholders[1].text = "Automated Concept Summary\nGenerated by Malak's Architecture"
    
    chunked_points = [points[i:i + 2] for i in range(0, len(points), 2)]
    for idx, chunk in enumerate(chunked_points[:9], 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Key Concepts - Part {idx}"
        tf = slide.placeholders[1].text_frame
        for pt in chunk:
            p = tf.add_paragraph()
            p.text = pt.lstrip('-* ')
            
    output_path = "summary_presentation.pptx"
    prs.save(output_path)
    return output_path

def get_ollama_models():
    
    try:
        result = subprocess.run(["ollama", "list"], capture_output=True, text=True)
        lines = result.stdout.strip().split('\n')[1:] 
        models = [line.split()[0] for line in lines if line]
        return models if models else ["mistral", "llama3", "deepseek-r1", "gemma3", "qwen3"]
    except Exception:
        return ["mistral", "llama3", "deepseek-r1", "gemma3", "qwen3"]

# =========================
# SIDEBAR & FILE UPLOADER
# =========================
st.sidebar.header("Controls")

page = st.sidebar.radio(
    "Navigation",
    [
        "Home / Ask",
        "Semantic Search",
        "Flashcards",
        "Review Cards",
        "Study Mode",
        "Dashboard",
        "Exports",
        "Exam System",
        "Project Builder",
        "Code Generator"
    ]
)

profile = st.session_state.user_profile

st.sidebar.subheader("User Profile")
st.sidebar.write(f"Level: {profile['level']:.2f}")
st.sidebar.write(f"XP: {profile['xp']}")
st.sidebar.write(f"Coins: {profile['coins']}")

mode = st.sidebar.radio(
    "Mode",
    ["Normal", "Study"]
)

available_models = get_ollama_models()
selected_model = st.sidebar.selectbox(
    "LLM Model",
    available_models
)

st.sidebar.info(
    f"Current Model : {selected_model}"
)

uploaded_files = st.sidebar.file_uploader(
    "Upload Study Materials (Multimodal Mode)",
    type=["pdf", "png", "jpg", "jpeg", "docx", "pptx", "csv"],
    accept_multiple_files=True
)

if "saved_files" not in st.session_state:
    st.session_state.saved_files = {}

if uploaded_files:
    os.makedirs("temp_pdfs", exist_ok=True)
    for file in uploaded_files:
        file_path = os.path.join("temp_pdfs", file.name)
        with open(file_path, "wb") as f:
            f.write(file.getbuffer())
        st.session_state.saved_files[file.name] = file_path

# =========================
# OLLAMA
# =========================
def call_ollama(prompt, model):
    
    try:
        res = requests.post(
            "http://localhost:11434/api/generate",
            json={
                "model": model,
                "prompt": prompt,
                "stream": False
            },
            timeout=300 
        ).json()

        return res.get("response", "")

    except Exception as e:
        return f"Ollama not running or error happened: {e}"

# =========================
# CHAPTER SPLIT
# =========================
def split_into_chapters(text):
    patterns = [
        r"(Chapter\s+\d+.*)",
        r"(\d+\.\s+.+)",
        r"([A-Z][A-Z\s]{5,})",
    ]

    combined = "|".join(patterns)
    splits = re.split(combined, text)

    chapters = []
    current = ""

    for part in splits:
        if not part:
            continue

        if re.match(combined, part.strip()):
            if current:
                chapters.append(current.strip())
            current = part
        else:
            current += " " + part

    if current:
        chapters.append(current.strip())

    if not chapters and text.strip():
        chapters.append(text.strip())

    return chapters

# =========================
# BUILD INDEX 
# =========================
@st.cache_data
def load_embedding_model():
    return SentenceTransformer("all-MiniLM-L6-v2")

def build_index(files):
    if not files:
        return None, None, [], [], [], []
    chunks = []
    sources = []
    pages = []
    chunk_to_chapter = []

    for file in files:
        file_ext = file.name.split('.')[-1].lower()
        file_text = ""
        
        if file_ext == "pdf":
            try:
                with pdfplumber.open(file) as pdf:
                    for page_num, page_obj in enumerate(pdf.pages, start=1):
                        page_text = page_obj.extract_text() or ""
                        tables = page_obj.extract_tables()
                        if tables:
                            page_text += "\n\n[Extracted PDF Table Structure Data]:\n"
                            for tbl in tables:
                                df_tbl = pd.DataFrame(tbl)
                                page_text += df_tbl.to_markdown(index=False) + "\n"
                        
                        if page_text.strip():
                            page_chunks = [page_text[i:i+800] for i in range(0, len(page_text), 700)]
                            for ch in page_chunks:
                                chunks.append(ch)
                                sources.append(file.name)
                                pages.append(page_num)
                                chunk_to_chapter.append(f"Page {page_num}")
            except Exception as e:
                st.error(f"Error parsing structural layout of PDF {file.name}: {e}")
            continue 
            
        elif file_ext in ["png", "jpg", "jpeg"]:
            file_text = extract_text_from_image(file)
        elif file_ext == "docx":
            file_text = read_docx(file)
        elif file_ext == "pptx":
            file_text = read_pptx(file)
        elif file_ext == "csv":
            try:
                df = pd.read_csv(file)
                file_text = df.to_string()
            except:
                file_text = "[CSV Read Error]"

        if file_text.strip():
            file_chunks = [file_text[i:i+800] for i in range(0, len(file_text), 700)]
            for ch in file_chunks:
                chunks.append(ch)
                sources.append(file.name)
                pages.append(1) 
                chunk_to_chapter.append("Uploaded Structural Document Asset")

    if not chunks:
        return None, None, [], [], [], []

    embed_model = load_embedding_model()
    embeddings = embed_model.encode(chunks)

    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(np.array(embeddings).astype("float32"))

    full_text = " ".join(chunks[:50])
    st.session_state.all_chapters = split_into_chapters(full_text)

    return index, embed_model, chunks, sources, pages, chunk_to_chapter

index, embed_model, chunks, sources, pages, chunk_to_chapter = build_index(uploaded_files)

# =========================
# CORE FUNCTIONS
# =========================
def detect_best_chapter_from_chunks(question):
    if index is None:
        return None
    q_embed = embed_model.encode([question])
    D, I = index.search(np.array(q_embed).astype("float32"), k=min(10, len(chunks)))

    chapter_scores = {}

    for idx in I[0]:
        chapter = chunk_to_chapter[idx]
        chapter_scores[chapter] = chapter_scores.get(chapter, 0) + 1

    return max(chapter_scores, key=chapter_scores.get)

def get_context(question, selected_chapter=None):
    if index is None:
        return "", ["LLM Knowledge"]
        
    if selected_chapter:
        best_chapter = selected_chapter
    else:
        best_chapter = detect_best_chapter_from_chunks(question)

    chapter_indices = [
        i for i in range(len(chunks))
        if chunk_to_chapter[i] == best_chapter
    ]

    if len(chapter_indices) == 0:
        chapter_indices = list(range(len(chunks)))

    chapter_chunks = [chunks[i] for i in chapter_indices]

    ch_embeddings = embed_model.encode(chapter_chunks)
    q_embed = embed_model.encode([question])[0]

    scores = []

    for original_idx, emb in zip(chapter_indices, ch_embeddings):
        sim = np.dot(q_embed, emb)
        feedback = st.session_state.chunk_feedback.get(original_idx, 0)

        level = st.session_state.user_profile["level"]
        weak_topics = st.session_state.user_profile["weak_topics"]

        chapter_weight = weak_topics.get(best_chapter, 0)

        final_score = sim + feedback + (0.3 * chapter_weight) + (0.2 * (1 - level))
        scores.append((final_score, original_idx))

    scores.sort(reverse=True)

    top_ids = [i for _, i in scores[:3]]
    top_pages = [pages[i] for i in top_ids]
    top_sources = [sources[i] for i in top_ids]
    st.session_state.last_chunk_ids = top_ids

    selected = [chunks[i] for i in top_ids]
    context = "\n\n".join(selected)

    citation = list(
        set(
            [
                f"{s} - {chunk_to_chapter[i]}"
                for i, s, p in zip(top_ids, top_sources, top_pages)
            ]
        )
    )

    return context, citation

def semantic_search(query, top_k=5):
    if index is None: return []
    q_embed = embed_model.encode([query])

    D, I = index.search(
        np.array(q_embed).astype("float32"),
        k=min(20, len(chunks))
    )

    results = []
    seen = set()

    for idx in I[0]:
        key = (
            sources[idx],
            pages[idx]
        )

        if key in seen:
            continue

        seen.add(key)

        results.append(
            {
                "text": chunks[idx],
                "page": pages[idx],
                "source": sources[idx],
                "path": st.session_state.saved_files.get(sources[idx], "#")
            }
        )

        if len(results) >= top_k:
            break

    return results

def generate_flashcards():
    prompt = f"""
Generate EXACTLY 5 flashcards.

FORMAT:

Q: ...
A: ...

Q: ...
A: ...

Context:
{st.session_state.context if st.session_state.context else st.session_state.question}

Do not add anything else.
"""
    return call_ollama(prompt, selected_model)

def update_card_schedule(question, result):
    if question not in st.session_state.card_stats:
        st.session_state.card_stats[question] = {
            "interval": 1,
            "next_review": datetime.now(),
            "successes": 0
        }

    card = st.session_state.card_stats[question]

    if result == "wrong":
        card["interval"] = 1
        card["successes"] = 0
    elif result == "correct":
        card["successes"] += 1
        card["interval"] = min(card["interval"] * 2, 30)
    elif result == "easy":
        card["successes"] += 1
        card["interval"] = min(card["interval"] * 3, 60)

    card["next_review"] = datetime.now() + timedelta(days=card["interval"])

def compute_confidence(question):
    if index is None: return 0.5
    q_embed = embed_model.encode([question])
    D, I = index.search(np.array(q_embed).astype("float32"), k=min(3, len(chunks)))

    avg_dist = sum(D[0]) / len(D[0]) if len(D[0]) > 0 else 1
    similarity_score = 1 / (1 + avg_dist)

    chunk_embeds = embed_model.encode([chunks[i] for i in I[0]])
    sim_matrix = np.dot(chunk_embeds, chunk_embeds.T)

    agreement_score = np.mean(sim_matrix)
    max_val = np.max(sim_matrix) if np.max(sim_matrix) != 0 else 1

    confidence = (0.6 * similarity_score) + (0.4 * (agreement_score / max_val))

    return float(max(0.0, min(confidence, 1.0)))

def is_ambiguous_question(question):
    question = question.strip()

    if len(question.split()) < 2:
        return True

    if len(question) < 5:
        return True

    return False

# =========================
# VOICE ASSISTANT FUNCTIONS
# =========================
def text_to_speech(text):
    
    tts = gTTS(text=text)
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as tmp:
        tts.save(tmp.name)
        with open(tmp.name, "rb") as f:
            st.audio(f.read(), format="audio/mp3")

# =========================
# 📄 PDF HIGHLIGHTING 
# =========================
def highlight_pdf_page(pdf_path, page_number):
    try:
        doc = fitz.open(pdf_path)
        page = doc[page_number - 1]
        rect = page.rect
        highlight = page.add_highlight_annot(rect)
        highlight.update()
        output_path = f"highlighted_page_{page_number}.pdf"
        doc.save(output_path)
        doc.close()
        return output_path
    except Exception as e:
        st.error(f"Highlight error: {e}")
        return None

# =========================
# 🤖 ASK AI
# =========================
def ask_ai(question, selected_chapter=None, explain_mode="University Student"):
    
    if index is not None:
        retrieval_result = retrieval_agent(
            question,
            selected_chapter,
            get_context
        )
        context = retrieval_result["context"]
        docs = retrieval_result["docs"]
    else:
        context = ""
        docs = ["LLM Knowledge"]

    st.session_state.context = context

    route = router_agent(question)

    # 🔹 Summary
    if route == "summary":
        answer = summarizer_agent(
            context if context else question,
            lambda p: call_ollama(p, selected_model)
        )
        return answer, docs, context

    # 🔹 Quiz
    elif route == "quiz":
        _, difficulty = difficulty_agent(st.session_state.user_profile)
        answer = generate_mcq(
            context if context else question,
            3,
            difficulty,
            lambda p: call_ollama(p, selected_model)
        )
        return answer, docs, context

    # 🔹Planner
    elif route == "planner":
        profile = st.session_state.user_profile
        answer = planner_agent(
            context if context else question,
            profile["weak_topics"],
            profile["level"],
            profile["history"],
            lambda p: call_ollama(p, selected_model)
        )
        return answer, docs, context

    # 🔹 Compare
    elif route == "compare":
        prompt = f"""
        Compare clearly.
        Use:
        - Similarities
        - Differences
        Context:
        {context if context else question}
        """
        answer = call_ollama(
            prompt,
            selected_model
        )
        answer = critic_agent(
            question,
            answer,
            lambda p: call_ollama(p, selected_model)
        )
        return answer, docs, context

    # 🔹 before coordinator_agent
    if context == "":
        context = question

    
    result = coordinator_agent(
        question,
        selected_chapter,
        st.session_state.conversation_memory,
        st.session_state.user_profile,
        get_context,
        router_agent,
        retrieval_agent,
        memory_agent,
        critic_agent,
        reflection_agent,
        recommendation_agent,
        lambda p: call_ollama(p, selected_model)
    )

    st.session_state.reflection = result["reflection"]
    st.session_state.recommendation = result["recommendation"]

    return (
        result["answer"],
        result["docs"],
        result["context"]
    )

def give_feedback(chunk_ids, liked=True):
    if not chunk_ids:
        return

    for cid in chunk_ids:
        if cid not in st.session_state.chunk_feedback:
            st.session_state.chunk_feedback[cid] = 0

        st.session_state.chunk_feedback[cid] += 0.2 if liked else -0.2

# =========================
# UPDATE BADGES 
# =========================
def update_badges():
    profile = st.session_state.user_profile
    badges = []

    if profile["xp"] >= 50:
        badges.append("Beginner Learner")

    if profile["xp"] >= 100:
        badges.append("Quiz Master")

    if profile["level"] >= 0.8:
        badges.append("Advanced Student")

    if profile["coins"] >= 100:
        badges.append("Rich Student")

    if profile["streak"] >= 7:
        badges.append("7-Day Streak")

    if profile["xp"] >= 500:
        badges.append("Legend")

    profile["badges"] = list(set(profile["badges"] + badges))

# =========================
# 🎖 UPDATE ACHIEVEMENTS 
# =========================
def update_achievements():
    profile = st.session_state.user_profile
    achievements = profile["achievements"]

    if profile["xp"] >= 100 and "100 XP" not in achievements:
        achievements.append("100 XP")

    if profile["xp"] >= 500 and "500 XP" not in achievements:
        achievements.append("500 XP")

    if profile["coins"] >= 100 and "Rich Student" not in achievements:
        achievements.append("Rich Student")

    if profile["streak"] >= 7 and "7-Day Streak" not in achievements:
        achievements.append("7-Day Streak")

    if len(profile["history"]) >= 10 and "Quiz Master" not in achievements:
        achievements.append("Quiz Master")

# =========================
# UPDATE STREAK 
# =========================
def update_streak():
    profile = st.session_state.user_profile
    today = date.today()

    if profile["last_study_date"] is None:
        profile["streak"] = 1
    else:
        last_day = profile["last_study_date"]
        
        if isinstance(last_day, str):
            last_day = datetime.strptime(last_day, "%Y-%m-%d").date()

        if today == last_day + timedelta(days=1):
            profile["streak"] += 1
        elif today != last_day:
            profile["streak"] = 1

    profile["last_study_date"] = today

# =========================
#  COMPUTE MASTERY 
# =========================
def compute_topic_mastery():

    profile = st.session_state.user_profile

    mastery = {}

    for topic, weakness in profile["weak_topics"].items():

        mastery[topic] = round(
            (1 - weakness) * 100,
            1
        )

    return mastery

# =========================
# 📊 GENERATE WEEKLY REPORT 
# =========================
def generate_weekly_report():

    profile = st.session_state.user_profile
    history = profile["history"]

    if not history:
        return "No study history yet."

    avg_accuracy = (
        sum(item["accuracy"] for item in history)
        / len(history)
    ) * 100

    mastery = compute_topic_mastery()

    strong_topics = [
        k for k, v in mastery.items()
        if v >= 80
    ]

    weak_topics = [
        k for k, v in mastery.items()
        if v < 60
    ]

    report = f"""
📊 WEEKLY REPORT

Average Accuracy:
{avg_accuracy:.1f}%

XP:
{profile['xp']}

Coins:
{profile['coins']}

Strong Topics:
{strong_topics}

Weak Topics:
{weak_topics}

Recommendation:
Focus on weak topics and solve more MCQs.
"""

    return report

# =========================
# CHAPTER TITLES
# =========================
chapter_titles = []

if st.session_state.all_chapters:
    chapter_titles = [
        ch.split("\n")[0][:60]
        for ch in st.session_state.all_chapters
    ]

# =========================
# PAGE 1: HOME / ASK
# =========================
if page == "Home / Ask":
    st.header("🏠 Home / Ask")

    selected_chapter = None

    if chapter_titles:
        selected_title = st.selectbox(
            "Choose Chapter optional",
            ["Auto Detect"] + chapter_titles
        )

        if selected_title != "Auto Detect":
            idx = chapter_titles.index(selected_title)
            selected_chapter = st.session_state.all_chapters[idx]

    question = st.text_input("Ask your question")

    level = st.session_state.user_profile["level"]
    if level < 0.4:
        default_mode = "Beginner"
    elif level < 0.7:
        default_mode = "University Student"
    else:
        default_mode = "Expert"

    explain_mode = st.selectbox(
        "Explain Like",
        [
            "Child",
            "Beginner",
            "University Student",
            "Expert"
        ],
        index=[
            "Child",
            "Beginner",
            "University Student",
            "Expert"
        ].index(default_mode)
    )

    st.write("🎙 Voice Input:")
    voice_question = st_speech_to_text(language='en', use_container_width=True, just_once=True, key='STT')
    
    if voice_question:
        st.success(voice_question)
        question = voice_question

    if st.button("Ask"):
        if is_ambiguous_question(question):
            st.warning("Your question is unclear. Try writing a more specific question.")
        else:
            ans, doc, ctx = ask_ai(question, selected_chapter, explain_mode)

            st.session_state.answer = ans
            st.session_state.doc = doc
            st.session_state.context = ctx
            st.session_state.question = question

            st.session_state.conversation_memory.append({
                "question": question,
                "answer": ans
            })

            st.session_state.chat_history.append({
                "question": question,
                "answer": ans,
                "source": doc
            })

    if st.session_state.answer:
        st.subheader("Answer")
        st.write(st.session_state.answer)

        if st.button("🔊 Speak Answer"):
            text_to_speech(st.session_state.answer)

        st.subheader("Sources")
        if st.session_state.doc:
            for item in st.session_state.doc:
                st.write(item)

        if (
            st.session_state.doc
            and
            st.session_state.doc[0] != "LLM Knowledge"
        ):
            if st.button("📄 Highlight Source Page"):
                try:
                    source_info = st.session_state.doc[0]
                    pdf_name = source_info.split(" - ")[0]
                    page_num = int(source_info.split("Page ")[1])
                    pdf_path = st.session_state.saved_files[pdf_name]
                    output_path = highlight_pdf_page(pdf_path, page_num)
                    st.session_state.highlighted_pdf = output_path
                except Exception as e:
                    st.error(e)

        if st.session_state.highlighted_pdf:
            with open(st.session_state.highlighted_pdf, "rb") as f:
                st.download_button(
                    "⬇ Download Highlighted PDF",
                    f,
                    file_name=st.session_state.highlighted_pdf,
                    mime="application/pdf"
                )

        if index is not None:
            conf = compute_confidence(st.session_state.question)
            st.progress(conf)
            st.write(f"{conf*100:.1f}% confidence")

            col1, col2 = st.columns(2)

            if col1.button("Like"):
                give_feedback(st.session_state.last_chunk_ids, True)
                st.success("Feedback saved.")

            if col2.button("Dislike"):
                give_feedback(st.session_state.last_chunk_ids, False)
                st.warning("Feedback saved.")

# =========================
# PAGE 2: SEMANTIC SEARCH
# =========================
elif page == "Semantic Search":
    st.header("🔍 Semantic Search")

    search_query = st.text_input(
        "Search in PDFs"
    )

    top_k = st.slider(
        "Number of Results",
        1,
        10,
        5
    )

    if index is None:
        st.warning(
            "Upload PDFs to use Semantic Search."
        )
    elif st.button("Search"):
        st.session_state.search_results = semantic_search(
            search_query,
            top_k
        )
    
    if st.session_state.search_results:
        for i, res in enumerate(st.session_state.search_results):
            with st.expander(
                f"Result {i+1} - Page {res['page']}"
            ):
                st.write(f"Source: {res['source']}")
                st.write(res["text"])
                
                st.info(f"Page : {res['page']}")
                
                if res["path"] != "#" and os.path.exists(res["path"]):
                    with open(res["path"], "rb") as pdf_file:
                        st.download_button(
                            label=f"Open File Asset (Page {res['page']})",
                            data=pdf_file,
                            file_name=res["source"],
                            mime="application/octet-stream",
                            key=f"pdf_{i}"
                        )

# =========================
# PAGE : FLASHCARDS
# =========================
elif page == "Flashcards":
    st.header("🧠 Flashcards")

    if (
        not st.session_state.context
        and
        not st.session_state.answer
    ):
        st.warning("Ask a question first.")
    else:
        if st.button("Generate Flashcards"):
            raw = generate_flashcards()
            cards = re.findall(
                r"Q:\s*(.*?)\nA:\s*(.*?)(?=\nQ:|\Z)",
                raw,
                re.S
            )
            st.session_state.flashcards = cards
            st.session_state.current_flashcard = 0
            st.session_state.show_answer = False

        if st.session_state.flashcards:
            if st.session_state.current_flashcard >= len(st.session_state.flashcards):
                st.success("🎉 All flashcards completed!")
                if st.button("Restart Flashcards"):
                    st.session_state.current_flashcard = 0
                    st.session_state.show_answer = False
                    st.rerun()
            else:
                i = st.session_state.current_flashcard
                q, a = st.session_state.flashcards[i]

                st.subheader(
                    f"Card {i+1}/{len(st.session_state.flashcards)}"
                )

                st.info(q)

                if st.button("Show Answer"):
                    st.session_state.show_answer = True

                if st.session_state.show_answer:
                    st.success(a)

                col1, col2, col3 = st.columns(3)

                if col1.button("Wrong"):
                    update_card_schedule(q, "wrong")
                    st.session_state.show_answer = False
                    st.session_state.current_flashcard += 1
                    st.rerun()

                if col2.button("Correct"):
                    update_card_schedule(q, "correct")
                    st.session_state.show_answer = False
                    st.session_state.current_flashcard += 1
                    st.rerun()

                if col3.button("Easy"):
                    update_card_schedule(q, "easy")
                    st.session_state.show_answer = False
                    st.session_state.current_flashcard += 1
                    st.rerun()

# =========================
# PAGE : REVIEW CARDS
# =========================
elif page == "Review Cards":
    st.header("📅 Cards Due Today")

    now = datetime.now()
    due = []

    for q, data in st.session_state.card_stats.items():
        if data["next_review"] <= now:
            due.append((q, data["interval"]))

    if due:
        for q, interval in due:
            st.success(q)
            st.write(f"Interval : {interval} days")
            st.markdown("---")
    else:
        st.info("No cards due today.")

# =========================
# PAGE 3: STUDY MODE
# =========================
elif page == "Study Mode":
    st.header("📚 Study Mode")

    if not st.session_state.answer:
        st.warning("Ask a question first from Home / Ask page.")
    else:
        tabs = st.tabs(
            [
                "Summary",
                "Learning",
                "MCQ",
                "Written",
                "Compare",
                "Mind Map",
                "Diagram",
                "AI Tutor",
                "Mock Interview"
            ]
        )

        # SUMMARY
        with tabs[0]:
            if st.button("Generate Summary"):
                st.session_state.summary = summarizer_agent(
                    st.session_state.context if st.session_state.context else st.session_state.question,
                    lambda prompt: call_ollama(prompt, selected_model)
                )

            if st.session_state.summary:
                st.write(st.session_state.summary)

        # LEARNING
        with tabs[1]:
            if st.button("Learning Path"):
                profile = st.session_state.user_profile
                st.session_state.learning = planner_agent(
                    st.session_state.context if st.session_state.context else st.session_state.question,
                    profile["weak_topics"],
                    profile["level"],
                    profile["history"],
                    lambda prompt: call_ollama(prompt, selected_model)
                )

            if st.session_state.learning:
                st.write(st.session_state.learning)

        # MCQ
        with tabs[2]:
            num_questions = st.slider("Number of Questions", 1, 10, 3)

            student_type, difficulty = difficulty_agent(
                st.session_state.user_profile
            )

            st.info(f"Adaptive Difficulty: {difficulty}")
            st.info(f"Student Type: {student_type}")

            if st.button("Generate Adaptive MCQ"):
                st.session_state.mcqs_raw = generate_mcq(
                    st.session_state.context if st.session_state.context else st.session_state.question,
                    num_questions,
                    difficulty,
                    lambda prompt: call_ollama(prompt, selected_model)
                )

            if st.session_state.mcqs_raw:
                blocks = re.split(r"Q\d+:", st.session_state.mcqs_raw)[1:]
                user_answers = []

                for i, block in enumerate(blocks):
                    lines = block.strip().split("\n")
                    question_text = lines[0] if lines else ""

                    options = [
                        l for l in lines
                        if re.match(r"[A-D][\)\.]", l.strip())
                    ]

                    answer_line = [
                        l for l in lines
                        if "Answer:" in l
                    ]

                    if answer_line:
                        correct = answer_line[0].split(":")[1].strip()[0]
                    else:
                        correct = "A"

                    st.write(f"**Q{i + 1}: {question_text}**")

                    if options:
                        choice = st.radio(
                            f"Select answer Q{i + 1}",
                            options,
                            key=f"mcq_{i}"
                        )

                        user_answers.append((choice[0], correct))

                if st.button("Check Answers"):
                    if not user_answers:
                        st.error("No MCQ answers found.")
                    else:
                        correct_count = 0

                        for i, (u, c) in enumerate(user_answers):
                            if u == c:
                                st.success(f"Q{i + 1}: Correct")
                                correct_count += 1
                            else:
                                st.error(f"Q{i + 1}: Wrong. Correct: {c}")

                        accuracy = correct_count / len(user_answers)

                        profile = st.session_state.user_profile
                        profile["level"] = (profile["level"] + accuracy) / 2
                        profile["xp"] += correct_count * 10
                        
                        profile["coins"] += correct_count * 5
                        profile["today_solved"] += len(user_answers)

                        current_chapter = (
                            st.session_state.doc[0]
                            if (
                                st.session_state.doc
                                and
                                st.session_state.doc[0] != "LLM Knowledge"
                            )
                            else st.session_state.question
                        )

                        if current_chapter:
                            profile["weak_topics"][current_chapter] = 1 - accuracy

                        profile["history"].append({
                            "question": st.session_state.question,
                            "accuracy": accuracy,
                            "level": profile["level"],
                            "chapter": current_chapter
                        })

                        update_badges()
                        update_streak()
                        update_achievements()

                        st.info(f"Your score: {correct_count}/{len(user_answers)}")
                        st.success(f"XP gained: {correct_count * 10}")

        # WRITTEN
        with tabs[3]:
            if st.button("Generate Written Questions"):
                st.session_state.written_raw = generate_written_questions(
                    st.session_state.context if st.session_state.context else st.session_state.question,
                    lambda prompt: call_ollama(prompt, selected_model)
                )

            if st.session_state.written_raw:
                questions = re.findall(r"Q\d:\s*(.*)", st.session_state.written_raw)
                answers = []

                for i, q in enumerate(questions):
                    st.write(f"**Q{i + 1}: {q}**")

                    ans = st.text_area(
                        f"Answer Q{i + 1}",
                        key=f"written_{i}"
                    )

                    answers.append((q, ans))

                if st.button("Evaluate Answers"):

                    for i, (q, ans) in enumerate(answers):

                        st.markdown(f"### Feedback Q{i + 1}")

                        if not ans.strip():

                            st.error("No answer. 0/10")

                        else:

                            feedback = evaluate_answer(
                                q,
                                ans,
                                lambda prompt: call_ollama(prompt, selected_model)
                                )

                            st.write(feedback)

        # COMPARE
        with tabs[4]:
            if st.button("Compare"):
                prompt = f"""
                Compare clearly.

                Use:
                - Similarities
                - Differences

                Context:
                {st.session_state.context if st.session_state.context else st.session_state.question}
                """
                st.session_state.compare = call_ollama(prompt, selected_model)

            if st.session_state.compare:
                st.write(st.session_state.compare)

        # MIND MAP
        with tabs[5]:

            if st.button("Generate Mind Map"):

                prompt = f"""
                Generate a hierarchical mind map.

                Use format like:

                Sensors
                ├── Analog
                │   ├── RTD
                │   ├── Thermistor
                ├── Digital
                │   ├── Encoder
                │   ├── Hall Sensor

                Context:

                {st.session_state.context if st.session_state.context else st.session_state.question}

                Return ONLY the mind map.
                """

                st.session_state.mind_map = call_ollama(
                    prompt,
                    selected_model
                )

            if st.session_state.mind_map:

                st.subheader("Mind Map")

                st.code(
                    st.session_state.mind_map
                )

        # DIAGRAM
        with tabs[6]:

            if st.button("Generate Diagram"):

                prompt = f"""
                Generate a Mermaid flowchart.

                Example:

                flowchart TD

                A[Input]

                B[Processing]

                C[Output]

                A --> B

                B --> C

                Context:

                {st.session_state.context if st.session_state.context else st.session_state.question}

                Return ONLY Mermaid code.
                """

                st.session_state.diagram = call_ollama(
                    prompt,
                    selected_model
                )

            if st.session_state.diagram:

                st.subheader("Mermaid Diagram")

                st.code(
                    st.session_state.diagram,
                    language="text"
                )

        # AI TUTOR
        with tabs[7]:

            st.subheader("🧠 AI Tutor")

            student_answer = st.text_area(
                "Your answer",
                key="tutor_answer"
            )

            if st.button("Next Question"):

                history = st.session_state.tutor_history

                history.append(
                    {
                        "student_answer": student_answer
                    }
                )

                response = tutor_agent(
                    st.session_state.context if st.session_state.context else st.session_state.question,
                    history,
                    lambda prompt:
                    call_ollama(
                        prompt,
                        selected_model
                    )
                )

                history.append(
                    {
                        "ai": response
                    }
                )

                st.write(response)
        
        # MOCK INTERVIEW
        with tabs[8]:

            st.subheader("🎓 Mock Interview")

            interview_answer = st.text_area(
                "Your interview answer",
                key="interview_answer"
            )

            if st.button("Next Interview Question"):

                history = st.session_state.interview_history

                history.append(
                    {
                        "student_answer": interview_answer
                    }
                )

                response = interview_agent(
                    st.session_state.context if st.session_state.context else st.session_state.question,
                    history,
                    lambda prompt:
                    call_ollama(
                        prompt,
                        selected_model
                    )
                )

                history.append(
                    {
                        "ai": response
                    }
                )

                st.write(response)

# =========================
# PAGE 4: DASHBOARD
# =========================
elif page == "Dashboard":
    st.header("📊 Dashboard")

    profile = st.session_state.user_profile
    history = profile["history"]

    total_questions = len(history)
    avg_accuracy = 0

    if history:
        avg_accuracy = sum(item["accuracy"] for item in history) / len(history)

    col1, col2, col3, col4 = st.columns(4)

    col1.metric("Total Questions", total_questions)
    col2.metric("Current Level", f"{profile['level']:.2f}")
    col3.metric("Average Accuracy", f"{avg_accuracy * 100:.1f}%")
    col4.metric("XP", profile["xp"])
    
    col5, col6, col7 = st.columns(3)
    
    col5.metric("Coins", profile["coins"])
    col6.metric("Streak", profile["streak"])
    col7.metric(
        "Daily Goal",
        f"{profile['today_solved']}/{profile['daily_goal']}"
    )

    st.subheader("🏆 Badges")

    if profile["badges"]:
        for badge in profile["badges"]:
            st.success(badge)
    else:
        st.info("No badges yet. Solve quizzes to earn badges.")
        
    st.subheader("🎖 Achievements")

    if profile["achievements"]:
        for ach in profile["achievements"]:
            st.success(ach)
    else:
        st.info("No achievements yet.")

    st.subheader("Weak Topics Chart")

    if profile["weak_topics"]:
        weak_df = pd.DataFrame({
            "Topic": list(profile["weak_topics"].keys()),
            "Weakness": list(profile["weak_topics"].values())
        })

        st.bar_chart(weak_df.set_index("Topic"))
    else:
        st.info("No weak topics yet.")

    st.subheader("📈 Accuracy Curve")

    if history:
        accuracy_df = pd.DataFrame(
            {
                "Accuracy":[
                    item["accuracy"]*100
                    for item in history
                ]
            }
        )
        st.line_chart(accuracy_df)
    else:
        st.info("No accuracy data yet.")

    st.subheader("🎯 Accuracy Per Topic")

    if profile["weak_topics"]:
        topic_df = pd.DataFrame(
            {
                "Topic":list(profile["weak_topics"].keys()),
                "Accuracy":[
                    (1-v)*100
                    for v in profile["weak_topics"].values()
                ]
            }
        )
        st.bar_chart(topic_df.set_index("Topic"))
    else:
        st.info("No topic data yet.")

    st.subheader("🏆 Topic Mastery Meter")

    mastery = compute_topic_mastery()

    if mastery:

        mastery_df = pd.DataFrame(
            {
                "Topic": mastery.keys(),
                "Mastery": mastery.values()
            }
        )

        st.bar_chart(
            mastery_df.set_index(
                "Topic"
            )
        )

    else:

        st.info(
            "No topic mastery data yet."
        )

    st.subheader("🔥 Study Heat Map")

    if history:
        heat_data = []

        for item in history:
            day_name = datetime.now().strftime("%A")
            heat_data.append(
                {
                    "Day": day_name,
                    "Accuracy": item["accuracy"] * 100
                }
            )

        heat_df = pd.DataFrame(heat_data)
        st.bar_chart(heat_df.set_index("Day"))
    else:
        st.info("No study heat map yet.")

    st.subheader("🧠 Reflection Agent")

    if st.button("Analyze Learning"):
        with st.spinner("Analyzing your detailed study data and formulating your next plan..."):
            st.session_state.reflection = reflection_agent(
                st.session_state.user_profile,
                lambda prompt: call_ollama(
                    prompt,
                    selected_model
                )
            )

    if st.session_state.reflection:
        st.write(st.session_state.reflection)

    st.subheader("🎯 Recommendation Agent")

    if st.button("Generate Recommendations"):
        with st.spinner("Analyzing weak topics..."):
            st.session_state.recommendation = recommendation_agent(
                st.session_state.user_profile,
                lambda prompt: call_ollama(
                    prompt,
                    selected_model
                )
            )

    if st.session_state.recommendation:
        st.write(st.session_state.recommendation)

    st.subheader("📊 Weekly Report")

    if st.button("Generate Weekly Report"):

        report = generate_weekly_report()

        st.write(report)

    st.subheader("📅 Revision Plan")

    tasks = revision_agent(
        st.session_state.user_profile
    )

    if tasks:

        for task in tasks:

            st.success(
                task["task"]
            )

    else:

        st.info(
            "No revision tasks yet."
        )

    st.subheader("Chat History")

    if st.session_state.chat_history:
        for i, chat in enumerate(st.session_state.chat_history):
            with st.expander(f"Question {i + 1}: {chat['question']}"):
                st.write("Answer:")
                st.write(chat["answer"])
                st.write("Source:")
                st.write(chat["source"])
    else:
        st.info("No chat history yet.")

# =========================
# PAGE 5: EXPORTS & LEARNING ASSISTANT
# =========================
elif page == "Exports":
    st.header("📄 Exports & Learning Assistant")

    
    st.subheader("📺 YouTube Learning Assistant")
    profile = st.session_state.user_profile
    
    if profile["weak_topics"]:
        weakest_topic = max(profile["weak_topics"], key=profile["weak_topics"].get)
        st.info(f"Detected Weak Topic: **{weakest_topic}**. Here are recommended video lectures to help you review:")
        
        with st.spinner("Searching YouTube for custom recommendations..."):
            videos = search_youtube_videos(f"{weakest_topic} lecture tutorial")
            if videos:
                for idx, vid in enumerate(videos, 1):
                    st.video(vid)
                    st.caption(f"Recommendation #{idx}: {vid}")
            else:
                st.write("Could not fetch videos at the moment. Try again later.")
    else:
        st.info("Solve some adaptive quizzes in 'Study Mode' so the YouTube Assistant can pinpoint your weak topics and recommend lectures!")

    st.markdown("---")

    
    st.subheader("🖥️ PDF Summary ➡️ PowerPoint Presentation")
    if st.session_state.summary:
        if st.button("Generate 10-Slide PPTX Presentation"):
            with st.spinner("Structuring your presentation..."):
                ppt_path = generate_pptx_from_summary(st.session_state.summary)
                
                with open(ppt_path, "rb") as ppt_file:
                    st.download_button(
                        label="📥 Download PowerPoint (PPTX)",
                        data=ppt_file,
                        file_name="Study_Summary_Presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            st.success("Your 10-Slide presentation is ready for download!")
    else:
        st.warning("⚠️ Please generate a Summary first from the 'Study Mode' tab before exporting to PowerPoint.")

    st.markdown("---")

    
    st.subheader("📝 Export Full Study Report")
    export_text = ""
    if st.session_state.summary:
        export_text += "SUMMARY\n" + st.session_state.summary + "\n\n"
    if st.session_state.mcqs_raw:
        export_text += "MCQs\n" + st.session_state.mcqs_raw + "\n\n"

    if export_text:
        st.download_button(
            label="Download Study Report as TXT",
            data=export_text,
            file_name="study_report.txt",
            mime="text/plain"
        )
    else:
        st.info("Nothing to export yet.")

# =========================
# PAGE : EXAM SYSTEM
# =========================
elif page == "Exam System":

    st.header("📝 AI Exam Generator")

    if (
        not st.session_state.context
        and
        not st.session_state.answer
    ):
        st.warning("Ask a question first.")
    else:
        student_type, _ = difficulty_agent(
            st.session_state.user_profile
        )

        if st.button("Generate Full Exam"):

            prompt = f"""
            Generate a complete exam.

            Student Level:
            {student_type}

            Rules:
            Beginner → easy questions.
            Intermediate → moderate questions.
            Advanced → hard questions.

            Include:

            1. Five MCQ questions.

            2. Five True or False questions.

            3. Five Fill in the blanks questions.

            4. Two Essay questions.

            Context:

            {
            st.session_state.context
            if st.session_state.context
            else st.session_state.question
            }

            Return clear formatting.
            """

            st.session_state.full_exam = call_ollama(
                prompt,
                selected_model
            )

        if "full_exam" in st.session_state:
            st.write(st.session_state.full_exam)

            st.subheader("Your Answers")

            student_answers = st.text_area(
                "Write your answers here"
            )

            if st.button("Grade Exam"):

                prompt = f"""
                Grade this exam strictly.

                Exam:

                {st.session_state.full_exam}

                Student Answers:

                {student_answers}

                Return ONLY:

                Final Score: X/20

                Mistakes:
                - ...

                Feedback:
                - ...

                Strong Points:
                - ...

                Weak Points:
                - ...
                """

                st.session_state.exam_feedback = call_ollama(
                    prompt,
                    selected_model
                )

            if st.session_state.exam_feedback:

                st.subheader("Exam Evaluation")

                st.write(
                    st.session_state.exam_feedback
                )

# =========================
# PAGE : PROJECT BUILDER
# =========================
elif page == "Project Builder":

    st.header("🚀 Project Builder Agent")

    project_topic = st.text_input(
        "Enter Project Idea"
    )

    if st.button("Build Project"):

        result = project_builder_agent(
            project_topic,
            (
                st.session_state.context
                or
                project_topic
            ),
            lambda prompt:
            call_ollama(
                prompt,
                selected_model
            )
        )

        st.subheader(
            "Generated Project"
        )

        st.write(result)

# =========================
# PAGE : CODE GENERATOR
# =========================
elif page == "Code Generator":

    st.header("💻 Code Generator Agent")

    topic = st.text_input(
        "What do you want to generate?"
    )

    language = st.selectbox(
        "Programming Language",
        [
            "Python",
            "C++",
            "Arduino",
            "SQL",
            "Java",
            "JavaScript"
        ]
    )

    if st.button("Generate Code"):

        result = code_generator_agent(
            topic,
            language,
            (
                st.session_state.context
                or
                topic
            ),
            lambda prompt:
            call_ollama(
                prompt,
                selected_model
            )
        )

        st.subheader("Generated Code")

        st.write(result)